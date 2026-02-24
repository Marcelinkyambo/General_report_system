import streamlit as st
import pandas as pd
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.util import Inches
import tempfile

# --------------------------------------------------
# Streamlit setup
# --------------------------------------------------
st.set_page_config(page_title="Sales Performance Report", layout="wide")
st.title("Sales Performance → PowerPoint Generator")

# --------------------------------------------------
# Helpers
# --------------------------------------------------
def save_plot(fig):
    tmp = tempfile.NamedTemporaryFile(delete=False, suffix=".png")
    fig.savefig(tmp.name, dpi=200, bbox_inches="tight")
    plt.close(fig)
    return tmp.name


def add_title_slide(prs, title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle


def add_chart_slide(prs, title, image):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = title
    slide.shapes.add_picture(image, Inches(1), Inches(1.5), width=Inches(8))


def add_text_slide(prs, title, text):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = title
    slide.shapes.add_textbox(
        Inches(1.5), Inches(3), Inches(7), Inches(1.5)
    ).text_frame.text = text


# --------------------------------------------------
# Upload
# --------------------------------------------------
file = st.file_uploader("Upload Sales Report (Excel or CSV)", type=["xlsx", "csv"])

if file:
    df = pd.read_excel(file) if file.name.endswith("xlsx") else pd.read_csv(file)

    # --------------------------------------------------
    # Validate structure
    # --------------------------------------------------
    required_cols = {
        "SKU",
        "Item Name",
        "Category",
        "Sales Value Current Year",
        "Sales Current Year  ",
        "Closing Balance"
    }

    if not required_cols.issubset(df.columns):
        st.error("Uploaded report structure does not match expected format.")
        st.write("Detected columns:", list(df.columns))
        st.stop()

    # --------------------------------------------------
    # Data cleaning (KEEP ZERO SALES)
    # --------------------------------------------------
    for col in ["Sales Value Current Year", "Sales Current Year  ", "Closing Balance"]:
        df[col] = pd.to_numeric(df[col], errors="coerce").fillna(0)

    df["Category"] = df["Category"].fillna("Uncategorized")
    df["Item Name"] = df["Item Name"].fillna("Unknown Item")
    df["Brand"] = df["Item Name"].astype(str).str.split().str[0]

    prs = Presentation()

    # --------------------------------------------------
    # SLIDE 1: COVER
    # --------------------------------------------------
    add_title_slide(
        prs,
        "Sales Performance Analysis",
        "Year to Date – December 2025"
    )

    # --------------------------------------------------
    # SLIDE 2: CATEGORY CONTRIBUTION (PIE)
    # --------------------------------------------------
    category_sales = df.groupby("Category")["Sales Value Current Year"].sum()
    category_chart = category_sales[category_sales > 0]

    if len(category_chart) > 0:
        fig, ax = plt.subplots()
        ax.pie(
            category_chart.values,
            labels=category_chart.index,
            autopct="%1.1f%%",
            startangle=90
        )
        ax.set_title("Top Categories by Total Sales Value")
        add_chart_slide(prs, "Category Contribution", save_plot(fig))
    else:
        add_text_slide(
            prs,
            "Category Contribution",
            "No sales recorded at category level for this period."
        )

    # --------------------------------------------------
    # SLIDE 3: TOP 10 BRANDS
    # --------------------------------------------------
    brand_sales = (
        df.groupby("Brand")["Sales Value Current Year"]
        .sum()
        .sort_values(ascending=False)
        .head(10)
    )

    if len(brand_sales) > 0:
        fig, ax = plt.subplots(figsize=(8, 5))
        brand_sales.sort_values().plot(kind="barh", ax=ax)
        ax.set_title("Top 10 Brands by Sales Value")
        ax.set_xlabel("Sales Value (RWF)")
        add_chart_slide(prs, "Top 10 Brands", save_plot(fig))
    else:
        add_text_slide(
            prs,
            "Top 10 Brands",
            "No brand-level sales data available."
        )

    # --------------------------------------------------
    # ITEM SUMMARY (FULL DATASET)
    # --------------------------------------------------
    item_summary = (
        df.groupby("Item Name")
        .agg(
            sales_value=("Sales Value Current Year", "sum"),
            units=("Sales Current Year  ", "sum"),
            closing_stock=("Closing Balance", "sum")
        )
    )

    # --------------------------------------------------
    # SLIDE 4: TOP 20 BEST PERFORMING ITEMS
    # --------------------------------------------------
    top_items = (
        item_summary[item_summary["sales_value"] > 0]
        .sort_values("sales_value", ascending=False)
        .head(20)
    )

    if len(top_items) > 0:
        fig, ax = plt.subplots(figsize=(8, 7))
        top_items["sales_value"].sort_values().plot(kind="barh", ax=ax)
        ax.set_title("Top 20 Best Performing Items")
        ax.set_xlabel("Sales Value (RWF)")
        add_chart_slide(prs, "Best Performing Items", save_plot(fig))
    else:
        add_text_slide(
            prs,
            "Best Performing Items",
            "No items recorded sales during this period."
        )

    # --------------------------------------------------
    # SLIDE 5: TOP 20 LEAST PERFORMING ITEMS
    # (ZERO SALES + HIGH STOCK INCLUDED)
    # --------------------------------------------------
    worst_items = (
        item_summary
        .sort_values(["sales_value", "closing_stock"], ascending=[True, False])
        .head(20)
    )

    fig, ax = plt.subplots(figsize=(8, 7))
    worst_items["closing_stock"].plot(kind="barh", ax=ax)
    ax.set_title("Top 20 Least Performing Items (High Stock, Low Sales)")
    ax.set_xlabel("Closing Stock")
    add_chart_slide(prs, "Least Performing Items", save_plot(fig))

    # --------------------------------------------------
    # EXPORT
    # --------------------------------------------------
    tmp = tempfile.NamedTemporaryFile(delete=False, suffix=".pptx")
    prs.save(tmp.name)

    with open(tmp.name, "rb") as f:
        st.download_button(
            "Download PowerPoint",
            f,
            file_name="Sales_Performance_Report_2025.pptx"
        )

    st.success("Presentation generated successfully.")
